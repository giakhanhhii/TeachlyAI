import base64
import json
import os
import re
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any


def _now_stamp() -> str:
    return time.strftime("%Y%m%d-%H%M%S")


def _safe_filename(name: str, default: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9._-]+", "_", name.strip())[:80].strip("_")
    return cleaned or default


def _json_loads_lenient(text: str) -> Any:
    """
    Parse JSON in a lenient way:
    - tries raw JSON
    - tries to extract the first {...} block
    """
    try:
        return json.loads(text)
    except Exception:
        pass

    m = re.search(r"\{.*\}", text, flags=re.DOTALL)
    if not m:
        raise ValueError("No JSON object found in model output")
    return json.loads(m.group(0))


class Chunker:
    """
    Agentic Chunking: use LLM to split an input document into coherent chunks.
    Fallback to heuristic splitting when LLM is unavailable.
    """

    def __init__(self, llm_client: Any | None, model: str):
        self._client = llm_client
        self._model = model

    def agentic_split(self, text: str, max_chunks: int = 12) -> list[str]:
        text = (text or "").strip()
        if not text:
            return []

        if not self._client:
            return self._fallback_split(text, max_chunks=max_chunks)

        prompt = (
            "Split the document into coherent, self-contained chunks for making lecture slides.\n"
            f"Rules:\n- Return ONLY valid JSON: {{\"chunks\": [\"...\"]}}\n"
            f"- Max chunks: {max_chunks}\n- Keep chunk order\n- Each chunk should be meaningful on its own\n\n"
            "DOCUMENT:\n"
            f"{text}"
        )

        try:
            raw = self._llm_json(prompt)
            chunks = raw.get("chunks")
            if not isinstance(chunks, list) or not all(isinstance(c, str) for c in chunks):
                raise ValueError("Invalid chunks format from LLM")
            chunks = [c.strip() for c in chunks if c.strip()]
            if not chunks:
                raise ValueError("LLM returned empty chunks")
            return chunks[:max_chunks]
        except Exception:
            return self._fallback_split(text, max_chunks=max_chunks)

    def _llm_json(self, prompt: str) -> dict[str, Any]:
        """
        Best-effort JSON generation across OpenAI SDK variants.
        """
        # Preferred (newer): responses API
        if hasattr(self._client, "responses"):
            resp = self._client.responses.create(
                model=self._model,
                input=[
                    {
                        "role": "user",
                        "content": [{"type": "input_text", "text": prompt}],
                    }
                ],
            )
            text = getattr(resp, "output_text", None) or ""
            return _json_loads_lenient(text)

        # Fallback: chat.completions
        resp = self._client.chat.completions.create(
            model=self._model,
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"},
        )
        msg = resp.choices[0].message
        return _json_loads_lenient(msg.content or "{}")

    @staticmethod
    def _fallback_split(text: str, max_chunks: int) -> list[str]:
        blocks = [b.strip() for b in re.split(r"\n\s*\n+", text) if b.strip()]
        if len(blocks) <= max_chunks:
            return blocks

        # If too many blocks, merge into roughly max_chunks groups.
        target = max(1, len(text) // max_chunks)
        chunks: list[str] = []
        buff: list[str] = []
        buff_len = 0
        for b in blocks:
            if buff and (buff_len + len(b) > target) and len(chunks) < max_chunks - 1:
                chunks.append("\n\n".join(buff).strip())
                buff = [b]
                buff_len = len(b)
            else:
                buff.append(b)
                buff_len += len(b)
        if buff:
            chunks.append("\n\n".join(buff).strip())
        return chunks[:max_chunks]


class SlideService:
    """
    Render a lecture script JSON into a PowerPoint file via python-pptx.
    """

    def __init__(self, output_dir: Path):
        self._output_dir = output_dir
        self._output_dir.mkdir(parents=True, exist_ok=True)

    def create(self, data: dict[str, Any], filename: str | None = None) -> Path:
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:  # pragma: no cover
            raise RuntimeError(
                "Missing dependency 'python-pptx'. Install it with: pip install python-pptx"
            ) from e

        title = str(data.get("title") or "Lecture").strip()
        topic = str(data.get("topic") or "").strip()
        slides = data.get("slides") or []
        if not isinstance(slides, list) or not slides:
            raise ValueError("Invalid lecture JSON: missing 'slides' array")

        if not filename:
            filename = f"{_safe_filename(title, 'lecture')}-{_now_stamp()}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        prs = Presentation()

        # Title slide
        layout0 = prs.slide_layouts[0]
        s0 = prs.slides.add_slide(layout0)
        s0.shapes.title.text = title
        if len(s0.placeholders) > 1:
            s0.placeholders[1].text = topic

        # Content slides
        for idx, slide in enumerate(slides, start=1):
            if not isinstance(slide, dict):
                continue
            slide_title = str(slide.get("title") or f"Slide {idx}").strip()
            content = str(slide.get("content") or "").strip()
            image_prompt = str(slide.get("image_prompt") or "").strip()

            layout = prs.slide_layouts[1]  # title + content
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide_title
            if len(s.placeholders) > 1:
                tf = s.placeholders[1].text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = content or "(No content)"
                if image_prompt:
                    tf.add_paragraph().text = f"[Image prompt] {image_prompt}"

        out_path = self._output_dir / filename
        prs.save(str(out_path))
        return out_path


@dataclass
class CircuitBreaker:
    failure_threshold: int = 2
    cooldown_seconds: int = 60

    _failures: int = 0
    _open_until: float = 0.0

    def allow(self) -> bool:
        return time.time() >= self._open_until

    def record_success(self) -> None:
        self._failures = 0
        self._open_until = 0.0

    def record_failure(self) -> None:
        self._failures += 1
        if self._failures >= self.failure_threshold:
            self._open_until = time.time() + self.cooldown_seconds


class VideoService:
    """
    Mock video generation pipeline:
    - TTS (mock) -> Image Gen (best-effort; may fail) -> FFmpeg (mock)
    Includes a basic circuit-breaker for the Image Gen step.
    """

    # A tiny valid 1x1 PNG (transparent)
    _PLACEHOLDER_PNG_B64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMA"
        "ASsJTYQAAAAASUVORK5CYII="
    )

    def __init__(self, output_dir: Path, llm_client: Any | None, image_model: str, enable_image_api: bool):
        self._output_dir = output_dir
        self._output_dir.mkdir(parents=True, exist_ok=True)
        self._client = llm_client
        self._image_model = image_model
        self._enable_image_api = enable_image_api
        self._cb = CircuitBreaker()

    def create(self, data: dict[str, Any], dir_name: str | None = None) -> Path:
        title = str(data.get("title") or "lecture").strip()
        slides = data.get("slides") or []
        if not isinstance(slides, list) or not slides:
            raise ValueError("Invalid lecture JSON: missing 'slides' array")

        if not dir_name:
            dir_name = f"{_safe_filename(title, 'lecture')}-video-{_now_stamp()}"
        out_dir = (self._output_dir / dir_name).resolve()
        out_dir.mkdir(parents=True, exist_ok=True)

        manifest: dict[str, Any] = {"title": title, "slides": []}

        for i, slide in enumerate(slides, start=1):
            if not isinstance(slide, dict):
                continue
            slide_title = str(slide.get("title") or f"Slide {i}").strip()
            content = str(slide.get("content") or "").strip()
            image_prompt = str(slide.get("image_prompt") or slide_title).strip()

            # 1) TTS mock
            narration_path = out_dir / f"{i:02d}-narration.txt"
            narration_path.write_text(f"{slide_title}\n\n{content}\n", encoding="utf-8")

            # 2) Image gen (best-effort with circuit breaker + fallback)
            image_path = out_dir / f"{i:02d}-image.png"
            used_fallback = False
            try:
                if not self._enable_image_api:
                    raise RuntimeError("Image API disabled (ENABLE_IMAGE_API=0)")
                if not self._client:
                    raise RuntimeError("No LLM client configured for image generation")
                if not self._cb.allow():
                    raise RuntimeError("Circuit breaker open; using placeholder image")
                self._generate_image(image_prompt, image_path)
                self._cb.record_success()
            except Exception:
                used_fallback = True
                self._cb.record_failure()
                self._write_placeholder_png(image_path)

            manifest["slides"].append(
                {
                    "slide_number": int(slide.get("slide_number") or i),
                    "title": slide_title,
                    "narration": str(narration_path.name),
                    "image": str(image_path.name),
                    "image_prompt": image_prompt,
                    "image_fallback": used_fallback,
                }
            )

        # 3) FFmpeg mock
        video_path = out_dir / "video.mock.txt"
        video_path.write_text(
            "Mock video pipeline complete.\n"
            "This is where you would stitch images + audio into an MP4 via FFmpeg.\n",
            encoding="utf-8",
        )

        (out_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        return out_dir

    def _generate_image(self, prompt: str, out_path: Path) -> None:
        """
        Best-effort OpenAI image generation.
        Falls back to placeholder outside this method.
        """
        if not hasattr(self._client, "images"):
            raise RuntimeError("OpenAI client does not support images API")

        resp = self._client.images.generate(
            model=self._image_model,
            prompt=prompt,
            size="1024x1024",
        )
        # OpenAI SDK returns base64 for images in some variants; handle both.
        b64 = None
        if getattr(resp, "data", None) and len(resp.data) > 0:
            item = resp.data[0]
            b64 = getattr(item, "b64_json", None) or getattr(item, "b64", None)
        if not b64:
            raise RuntimeError("Image API response did not include base64 data")
        out_path.write_bytes(base64.b64decode(b64))

    def _write_placeholder_png(self, out_path: Path) -> None:
        out_path.write_bytes(base64.b64decode(self._PLACEHOLDER_PNG_B64))

